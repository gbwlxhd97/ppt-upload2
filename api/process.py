from fastapi import FastAPI, File, UploadFile, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
import io

app = FastAPI()

# CORS 설정
origins = [
    "http://localhost:5173",
    "http://localhost:3000",
    "https://*.vercel.app",
    "*"
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.post("/api/slide-info")
async def get_slide_info(file: UploadFile = File(...)):
    """PPT 파일에서 슬라이드 개수 정보를 반환합니다."""
    try:
        pptx_content = await file.read()
        prs = Presentation(io.BytesIO(pptx_content))
        
        slide_count = len(prs.slides)
        
        return {
            "slide_count": slide_count,
            "filename": file.filename
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error getting slide info: {e}")

@app.post("/api/get-slide-text")
async def get_slide_text(
    file: UploadFile = File(...), 
    slide_index: int = Form(...)
):
    """슬라이드의 텍스트 정보만 반환 (클라이언트 렌더링용)"""
    try:
        contents = await file.read()
        prs = Presentation(io.BytesIO(contents))
        
        if not (0 <= slide_index < len(prs.slides)):
            raise HTTPException(status_code=400, detail="Invalid slide index")
            
        slide = prs.slides[slide_index]
        text_blocks = []
        
        # 각 텍스트 박스의 정보를 수집
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_blocks.append({
                    "text": shape.text,
                    "font_size": 16,  # 기본 폰트 크기
                    "bold": False
                })
        
        return {
            "slide_index": slide_index,
            "text_blocks": text_blocks,
            "full_text": "\n".join([block["text"] for block in text_blocks]) if text_blocks else f"Slide {slide_index + 1}"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error extracting slide text: {e}")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)